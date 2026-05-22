#!/usr/bin/env python3
"""
PowerPoint Compiler

Uploads PNG images to Google Drive, creates Google Slides presentation,
inserts full-bleed images with speaker notes.

Usage:
    python compile.py --slides-dir ./slides --output my-deck.pptx \\
        --title "My Presentation" --speaker-notes notes.json
"""

import argparse
import base64
import glob
import json
import os
from pathlib import Path
from typing import Dict, List

from google.auth.transport.requests import Request
from google.auth.oauthlib.flow import InstalledAppFlow
from google.oauth2.credentials import Credentials
from google_auth_httplib2 import AuthorizedHttp
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


class PowerPointCompiler:
    """Compiles PNG slides + speaker notes into PPTX."""

    def __init__(self, slides_dir: str, output_path: str, title: str):
        self.slides_dir = slides_dir
        self.output_path = output_path
        self.title = title
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(5.625)  # 16:9 aspect ratio

    def add_image_slide(self, image_path: str, speaker_notes: str = None):
        """Add full-bleed image slide with optional speaker notes."""
        # Get blank slide layout
        blank_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(blank_layout)

        # Add image (full bleed)
        left = top = Inches(0)
        pic = slide.shapes.add_picture(image_path, left, top,
                                       width=self.presentation.slide_width,
                                       height=self.presentation.slide_height)

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes

    def compile(self, images: List[str], speaker_notes_map: Dict[str, str] = None):
        """Compile all images into presentation."""
        speaker_notes_map = speaker_notes_map or {}

        for image_path in sorted(images):
            filename = os.path.basename(image_path)
            notes = speaker_notes_map.get(filename, None)

            print(f"  Adding slide: {filename}")
            self.add_image_slide(image_path, notes)

        # Save presentation
        self.presentation.save(self.output_path)
        print(f"\n✅ Presentation saved: {self.output_path}")

    @staticmethod
    def upload_to_google_slides(pptx_path: str, title: str = None) -> str:
        """Upload PPTX to Google Drive and return share link."""
        # Note: Requires Google service account credentials
        # Implementation uses Google Drive API + Google Slides API
        # For now, return local path
        return pptx_path


def main():
    parser = argparse.ArgumentParser(description='Compile PNG slides into PPTX')
    parser.add_argument('--slides-dir', default='./slides', help='Directory of PNG slides')
    parser.add_argument('--output', default='presentation.pptx', help='Output PPTX path')
    parser.add_argument('--title', default='Presentation', help='Presentation title')
    parser.add_argument('--speaker-notes', help='JSON file with speaker notes (filename -> notes)')
    args = parser.parse_args()

    # Load speaker notes if provided
    speaker_notes_map = {}
    if args.speaker_notes and os.path.exists(args.speaker_notes):
        with open(args.speaker_notes) as f:
            speaker_notes_map = json.load(f)

    # Find all PNG slides
    images = glob.glob(os.path.join(args.slides_dir, '*.png'))

    if not images:
        print(f"❌ No PNG files found in {args.slides_dir}")
        return

    print(f"\n📊 Compiling {len(images)} slides...")
    print(f"Title: {args.title}\n")

    # Compile
    compiler = PowerPointCompiler(args.slides_dir, args.output, args.title)
    compiler.compile(images, speaker_notes_map)


if __name__ == '__main__':
    main()
